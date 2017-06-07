from reportlab.pdfgen import canvas
from pptx import Presentation
import xlsxwriter
from docx import Document
from docx.shared import Inches
from ebooklib import epub
from PIL import Image,ImageFont,ImageDraw
from boto3.session import Session as boto3_session

def uploadToS3(strFolder,strFile,awsCred):
	if (awsCred['accessKeyId'] != ''):
		print('Uploading '+strFile+' to bucket '+awsCred['bucket']+'/'+strFolder)
		session = boto3_session(
					aws_access_key_id=awsCred['accessKeyId'],
					aws_secret_access_key=awsCred['secretAccessKey'],
					region_name="us-east-1")
		s3 = session.resource('s3')
		file_handle = open(strFile, 'rb')
		s3.Bucket(awsCred['bucket']).upload_file(file_handle.name, strFolder+'/'+strFile.split('/')[-1])
	else:
		print('Empty credentials, cannot upload')

def handler(event, context):

	awsCred = {}
	awsCred['accessKeyId'] = ''
	awsCred['secretAccessKey'] = ''
	awsCred['bucket'] = ''
	strFolder = 'docs'

	strFolderPath = '/tmp/'

	img = Image.new("RGB", (128, 128), "white")
	draw = ImageDraw.Draw(img)
	draw.text((10,10),"Hello world!",(0,0,0))
	img.save(strFolderPath+'demo.jpg')
	uploadToS3(strFolder,strFolderPath+'demo.jpg',awsCred)

	c = canvas.Canvas(strFolderPath+"demo.pdf")
	c.drawString(100,750,"Hello world!")
	c.save()
	uploadToS3(strFolder,strFolderPath+'demo.pdf',awsCred)

	document = Document()
	document.add_heading('Hello world!', level=1)
	document.add_page_break()
	document.save(strFolderPath+'demo.docx')
	uploadToS3(strFolder,strFolderPath+'demo.docx',awsCred)

	workbook = xlsxwriter.Workbook(strFolderPath+'demo.xlsx')
	worksheet = workbook.add_worksheet()
	worksheet.write('A1', 'Hello')
	worksheet.write('A2', 'World!')
	workbook.close()
	uploadToS3(strFolder,strFolderPath+'demo.xlsx',awsCred)

	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	title.text = "Hello, World!"
	prs.save(strFolderPath+'demo.pptx')
	uploadToS3(strFolder,strFolderPath+'demo.pptx',awsCred)

	book = epub.EpubBook()
	book.set_identifier('id123456')
	book.set_title('Hello world!')
	book.set_language('en')
	book.add_author('Author Authorowski')
	c1 = epub.EpubHtml(title='Hello world', file_name='chap_01.xhtml', lang='hr')
	c1.content=u'<h1>Hello world</h1><p>Hello world.</p>'
	book.add_item(c1)
	book.toc = (epub.Link('chap_01.xhtml', 'Introduction', 'intro'),
	             (epub.Section('Simple book'),
	             (c1, ))
	            )
	epub.write_epub(strFolderPath+'demo.epub', book, {})
	uploadToS3(strFolder,strFolderPath+'demo.epub',awsCred)

